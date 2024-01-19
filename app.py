# import sys 
# from langchain.embeddings import OpenAIEmbeddings
# from langchain.chat_models import ChatOpenAI
# from langchain.vectorstores import FAISS
from langchain_openai.chat_models import ChatOpenAI
from langchain_community.vectorstores import FAISS 
from langchain_openai.embeddings import OpenAIEmbeddings

import streamlit as st
import random
from dotenv import load_dotenv
import pypdfium2 as pdfium
from langchain.text_splitter import CharacterTextSplitter
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
import os
import pandas as pd
from pptx import Presentation
# from langchain.prompts import PromptTemplate
# from langchain.schema.messages import SystemMessage
# from sentence_transformers import SentenceTransformer
# from langchain.embeddings import HuggingFaceEmbeddings 

# print(sys.path)

# custom_template = """Given the following conversation and I want you to summarise the patient disease, and then identify the ICD Code that relates to this patient disease. And show me the patient procedure based on the ICD Code. If you do not know the answer reply with 'I am sorry'.
#     Chat History:
#     {chat_history}
#     Follow Up Input: {user_question}
#     """


def get_pdf_text(pdf_docs):
    text = ""
    for doc_index, pdf in enumerate(pdf_docs):
        # print(f"\033[93mReading PDF Document {doc_index + 1}/{len(pdf_docs)}\033[0m")
        pdf_reader = pdfium.PdfDocument(pdf) 
        num_pages = len(pdf_reader)
        # print(f"\033[94mNumber of pages in document: {num_pages}\033[0m")

        for i in range(num_pages):
            # print(f"\033[92mProcessing Page {i + 1}/{num_pages}...\033[0m")
            page = pdf_reader.get_page(i)
            textpage = page.get_textpage()
            page_text = textpage.get_text_range()
            text += page_text + "\n"

    # print("\033[91mPDF_TEXT IS =\033[0m", text)
    return text

def csv_to_pd(documents):
    #temp location for csv
    os.chdir('/temp-data')
    csv_files = [f for f in os.listdir() if f.endswith('.csv')]
    dfs = []
    for csv in documents:
        df = pd.read_csv(csv)
        dfs.append(df)
    combined_df = pd.concat(dfs, ignore_index=True)
    return combined_df

def pptx_to_text(documents):
    text = ''
    for eachfile in documents:
        prs = Presentation(eachfile)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    return text

def all_files(documents):
    text = ''
    for eachfile in documents:
        if eachfile.name.endswith('.csv'):
            text += csv_to_pd([eachfile])
        elif eachfile.name.endswith('.pdf'):
            text += get_pdf_text([eachfile])
        elif eachfile.name.endswith('.pptx'):
            text += pptx_to_text([eachfile])
    return(text)

#uses langchain function to split text -> so that it can be used for embedding
def get_text_chunks(documents):
    text_splitter = CharacterTextSplitter(separator="\n", chunk_size=5000, chunk_overlap=500, length_function=len)
    chunks = text_splitter.split_text(documents)
    return chunks

def get_vectorstore(text_chunks):
    # embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    embeddings = OpenAIEmbeddings() 
    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl") #can switch to instructor if stronger GPU, but im broke 
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    # index_name = init_pinecone()
    # vectorstore = Pinecone.from_documents(text_chunks, embedding= embeddings, index_name = index_name)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(temperature=0, model="gpt-4")
    # llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


def save_question_and_clear_prompt(ss):
    ss.user_question = ss.prompt_bar
    ss.prompt_bar = ""  # clearing the prompt bar after clicking enter to prevent automatic re-submissions


def write_chat(msgs):
    # Write the Q&A in a pretty chat format
    for i, msg in enumerate(msgs):
        if i % 2 == 0:  # it's a question
            st.markdown(user_template.replace("{{MSG}}", msg.content), unsafe_allow_html=True)
        else:  # it's an answer
            st.markdown(bot_template.replace("{{MSG}}", msg.content), unsafe_allow_html=True)



def main():
    load_dotenv()
    ss = st.session_state
    st.set_page_config(page_title="Chat with multiple PDFs",
                       layout="wide")
    st.write(css, unsafe_allow_html=True)
    st.markdown(css, unsafe_allow_html=True)
    st.image("https://1000logos.net/wp-content/uploads/2021/04/Accenture-logo-500x281.png", width=100)
    st.header("Ask questions about your documents")
    
    if "conversation_chain" not in ss:
        ss.conversation_chain = None  # the main variable storing the llm, retriever and memory
    if "prompt_bar" not in ss:
        ss.prompt_bar = ""
    if "user_question" not in ss:
        ss.user_question = ""
    if "docs_are_processed" not in ss:
        ss.docs_are_processed = False

    with st.sidebar:
        pdf_docs = st.file_uploader("Upload your PDFs here and click 'Process'", accept_multiple_files=True, type="pdf")
        if st.button("Process") and pdf_docs:
            with st.spinner("Processing"):
                raw_text = get_pdf_text(pdf_docs)  # get pdf text
                text_chunks = get_text_chunks(raw_text)  # get the text chunks
                vectorstore = get_vectorstore(text_chunks)  # create vector store
                ss.conversation_chain = get_conversation_chain(vectorstore)  # create conversation chain
                ss.docs_are_processed = True
            if ss.docs_are_processed:
                st.success("Files finished processing", icon="✅")
    st.text_input("Ask a question here:", key='prompt_bar', on_change=save_question_and_clear_prompt(ss))

    if ss.user_question:
        ss.conversation_chain({'question': ss.user_question})  # This is what gets the response from the LLM!
        if hasattr(ss.conversation_chain.memory, 'chat_memory'):
            chat = ss.conversation_chain.memory.chat_memory.messages
            write_chat(chat)

    if hasattr(ss.conversation_chain, 'memory'):  # There is memory if the documents have been processed
        if hasattr(ss.conversation_chain.memory, 'chat_memory'):  # There is chat_memory if questions have been asked
            if st.button("Forget conversation"):  # adding a button
                ss.conversation_chain.memory.chat_memory.clear()  # clears the ConversationBufferMemory

    # st.write(ss)  # use this when debugging for visualizing the session_state variables        


if __name__ == '__main__':
    main()

